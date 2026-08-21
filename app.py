"""
report_export.py
-----------------
Turns the AI narrative (Markdown from Groq) + the actual computed
statistics (DataFrames/dicts, not flattened text) into professional
DOCX, PDF, and PPTX files, each returned as raw bytes ready for
st.download_button.

All three build genuine tables for the statistics (not a monospace dump
of the LLM prompt text), a clean title page, and real bold formatting
for **markdown emphasis** in the narrative - instead of showing literal
'#' / '**' characters.

Kept dependency-light on purpose:
- DOCX: python-docx
- PDF:  reportlab (pure Python, no system binaries - installs cleanly
        on Streamlit Cloud, unlike weasyprint/wkhtmltopdf)
- PPTX: python-pptx
"""

import io
import re
from datetime import datetime

from docx import Document
from docx.shared import Pt, Inches as DocxInches, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors as rl_colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
    ListFlowable, ListItem, Table, TableStyle, HRFlowable,
)

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

# Shared "brand" color used across all three formats for a consistent look
ACCENT_HEX = "1F4E79"        # navy blue
ACCENT_RL = rl_colors.HexColor(f"#{ACCENT_HEX}")
ACCENT_DOCX = DocxRGBColor(0x1F, 0x4E, 0x79)
ACCENT_PPTX = PptxRGBColor(0x1F, 0x4E, 0x79)
LIGHT_BAND_RL = rl_colors.HexColor("#EEF3F8")


# --------------------------------------------------------------------------
# Small shared helpers
# --------------------------------------------------------------------------
def clean_dataset_title(filename: str) -> str:
    """'hotel_dataset (5).xlsx' -> 'Hotel Dataset'"""
    name = re.sub(r"\.[^.]+$", "", filename)          # drop extension
    name = re.sub(r"\s*\(\d+\)\s*$", "", name)          # drop ' (5)' duplicate suffix
    name = re.sub(r"[_\-]+", " ", name).strip()
    return name.title() if name else "Dataset"


def format_generated_at(dt: datetime) -> str:
    return dt.strftime("%d %B %Y, %H:%M")


def _parse_markdown_sections(markdown_text: str):
    """Split into [(heading_or_None, [raw_lines]), ...] on '## ' headings."""
    sections = []
    heading, lines = None, []
    for line in markdown_text.splitlines():
        if line.startswith("## "):
            if lines or heading:
                sections.append((heading, lines))
            heading, lines = line[3:].strip(), []
        else:
            lines.append(line)
    if lines or heading:
        sections.append((heading, lines))
    return sections


def _blocks_from_lines(lines):
    """Group section lines into ('bullets', [items]) / ('para', text) blocks."""
    blocks, bullet_buf, para_buf = [], [], []

    def flush():
        if bullet_buf:
            blocks.append(("bullets", bullet_buf[:]))
            bullet_buf.clear()
        if para_buf:
            blocks.append(("para", " ".join(para_buf)))
            para_buf.clear()

    for line in lines:
        stripped = line.strip()
        if not stripped:
            flush()
        elif stripped.startswith(("- ", "* ")):
            if para_buf:
                blocks.append(("para", " ".join(para_buf)))
                para_buf.clear()
            bullet_buf.append(stripped[2:].strip())
        else:
            if bullet_buf:
                blocks.append(("bullets", bullet_buf[:]))
                bullet_buf.clear()
            para_buf.append(stripped)
    flush()
    return blocks


_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def _split_bold_segments(text: str):
    """'**A** and B' -> [('A', True), (' and B', False)]"""
    segments, pos = [], 0
    for m in _BOLD_RE.finditer(text):
        if m.start() > pos:
            segments.append((text[pos:m.start()], False))
        segments.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        segments.append((text[pos:], False))
    return segments or [(text, False)]


def _markdown_to_reportlab_markup(text: str) -> str:
    """Escape XML specials, then turn **bold** into <b>bold</b> (ReportLab's
    Paragraph accepts a small HTML-like markup)."""
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return _BOLD_RE.sub(r"<b>\1</b>", text)


def _fmt_num(value, is_count=False) -> str:
    try:
        f = float(value)
    except (TypeError, ValueError):
        return str(value)
    if is_count:
        return f"{f:,.0f}"
    return f"{f:,.2f}"


def _render_fig_png(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    buf.seek(0)
    return buf


# --------------------------------------------------------------------------
# Shared: build plain "overview metrics" rows used by all three formats
# --------------------------------------------------------------------------
def _overview_rows(overview: dict):
    return [
        ("Rows", f"{overview['rows']:,}"),
        ("Columns", f"{overview['columns']:,}"),
        ("Duplicate rows", f"{overview['duplicate_rows']:,}"),
        ("Numeric columns", str(len(overview["numeric_columns"])) or "0"),
        ("Categorical columns", str(len(overview["categorical_columns"])) or "0"),
    ]


# ==========================================================================
# PDF (reportlab)
# ==========================================================================
def build_pdf_report(dataset_title, generated_at, model_name, narrative_md,
                      overview, missing_df, numeric_summary_df,
                      categorical_summary, corr_df, outliers, chart_fig=None) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER,
                             topMargin=0.7 * inch, bottomMargin=0.7 * inch,
                             leftMargin=0.7 * inch, rightMargin=0.7 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("title", parent=styles["Title"], textColor=ACCENT_RL, fontSize=24, spaceAfter=4)
    subtitle_style = ParagraphStyle("subtitle", parent=styles["Normal"], textColor=rl_colors.grey, fontSize=10, spaceAfter=16)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], textColor=ACCENT_RL, spaceBefore=16, spaceAfter=8)
    body = ParagraphStyle("body", parent=styles["Normal"], fontSize=10, leading=14)

    def table_style(header_bg=ACCENT_RL):
        return TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), header_bg),
            ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8.5),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [rl_colors.white, LIGHT_BAND_RL]),
            ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.HexColor("#CBD5E1")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ])

    story = [
        Paragraph(f"{dataset_title} — Exploratory Data Analysis Report", title_style),
        Paragraph(f"Generated {generated_at}  ·  Powered by Groq ({model_name})", subtitle_style),
        HRFlowable(width="100%", thickness=1, color=ACCENT_RL, spaceAfter=14),
    ]

    # Overview table
    story.append(Paragraph("Dataset Overview", h2))
    ov_rows = [["Metric", "Value"]] + [[k, v] for k, v in _overview_rows(overview)]
    t = Table(ov_rows, colWidths=[2.5 * inch, 3 * inch])
    t.setStyle(table_style())
    story.append(t)

    # AI narrative
    story.append(Paragraph("AI-Generated Insights", h2))
    for heading, lines in _parse_markdown_sections(narrative_md):
        if heading:
            story.append(Paragraph(heading, ParagraphStyle("h3", parent=styles["Heading3"], textColor=rl_colors.HexColor("#334155"))))
        for kind, content in _blocks_from_lines(lines):
            if kind == "bullets":
                items = [ListItem(Paragraph(_markdown_to_reportlab_markup(i), body)) for i in content]
                story.append(ListFlowable(items, bulletType="bullet", leftIndent=16, bulletColor=ACCENT_RL))
            else:
                story.append(Paragraph(_markdown_to_reportlab_markup(content), body))
        story.append(Spacer(1, 6))

    # Missing values
    story.append(Paragraph("Missing Values", h2))
    if missing_df is None or missing_df.empty:
        story.append(Paragraph("No missing values were detected in any column.", body))
    else:
        rows = [["Column", "Missing Count", "Missing %"]]
        for col, r in missing_df.iterrows():
            rows.append([str(col), f"{int(r['missing_count']):,}", f"{r['missing_pct']}%"])
        t = Table(rows, colWidths=[2.5 * inch, 1.5 * inch, 1.5 * inch])
        t.setStyle(table_style())
        story.append(t)

    # Numeric statistics
    if numeric_summary_df is not None and not numeric_summary_df.empty:
        story.append(Paragraph("Numeric Column Statistics", h2))
        cols = list(numeric_summary_df.columns)
        rows = [["Column"] + cols]
        for col, r in numeric_summary_df.iterrows():
            row = [str(col)] + [_fmt_num(r[c], is_count=(c == "count")) for c in cols]
            rows.append(row)
        t = Table(rows, colWidths=[1.3 * inch] + [0.72 * inch] * len(cols))
        t.setStyle(table_style())
        story.append(t)

    # Categorical summary
    if categorical_summary:
        story.append(Paragraph("Categorical Column Summary", h2))
        rows = [["Column", "Unique Values", "Top Values"]]
        for col, info in categorical_summary.items():
            top = ", ".join(f"{k} ({v})" for k, v in list(info["top_values"].items())[:3])
            rows.append([col, str(info["unique_values"]), top])
        t = Table(rows, colWidths=[1.6 * inch, 1.1 * inch, 3.3 * inch])
        t.setStyle(table_style())
        story.append(t)

    # Correlation
    if corr_df is not None and not corr_df.empty:
        story.append(Paragraph("Correlation Matrix", h2))
        cols = list(corr_df.columns)
        rows = [[""] + cols]
        for col in corr_df.index:
            rows.append([col] + [f"{corr_df.loc[col, c]:.2f}" for c in cols])
        col_w = min(0.85 * inch, 6 * inch / (len(cols) + 1))
        t = Table(rows, colWidths=[1.1 * inch] + [col_w] * len(cols))
        t.setStyle(table_style())
        story.append(t)

        if chart_fig is not None:
            story.append(Spacer(1, 10))
            story.append(RLImage(_render_fig_png(chart_fig), width=5.5 * inch, height=3.85 * inch, kind="proportional"))

    # Outliers
    if outliers:
        story.append(Paragraph("Potential Outliers (IQR rule)", h2))
        rows = [["Column", "Count", "% of values"]]
        for col, info in outliers.items():
            rows.append([col, str(info["count"]), f"{info['pct']}%"])
        t = Table(rows, colWidths=[2.5 * inch, 1.5 * inch, 1.5 * inch])
        t.setStyle(table_style())
        story.append(t)

    doc.build(story)
    return buf.getvalue()


# ==========================================================================
# DOCX
# ==========================================================================
def _docx_add_runs(paragraph, text, base_color=None):
    for segment, is_bold in _split_bold_segments(text):
        run = paragraph.add_run(segment)
        run.bold = is_bold
        if base_color:
            run.font.color.rgb = base_color


def _docx_table(doc, headers, rows, widths=None):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Light Grid Accent 1"
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = ""
        run = cell.paragraphs[0].add_run(h)
        run.bold = True
    for row in rows:
        cells = table.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = str(val)
    return table


def build_docx_report(dataset_title, generated_at, model_name, narrative_md,
                       overview, missing_df, numeric_summary_df,
                       categorical_summary, corr_df, outliers, chart_fig=None) -> bytes:
    doc = Document()

    title_p = doc.add_heading(level=0)
    title_run = title_p.add_run(f"{dataset_title} — Exploratory Data Analysis Report")
    title_run.font.color.rgb = ACCENT_DOCX

    meta_p = doc.add_paragraph()
    meta_run = meta_p.add_run(f"Generated {generated_at}  ·  Powered by Groq ({model_name})")
    meta_run.italic = True
    meta_run.font.size = Pt(9)
    meta_run.font.color.rgb = DocxRGBColor(0x64, 0x74, 0x8B)

    doc.add_heading("Dataset Overview", level=1)
    _docx_table(doc, ["Metric", "Value"], _overview_rows(overview))

    doc.add_heading("AI-Generated Insights", level=1)
    for heading, lines in _parse_markdown_sections(narrative_md):
        if heading:
            doc.add_heading(heading, level=2)
        for kind, content in _blocks_from_lines(lines):
            if kind == "bullets":
                for item in content:
                    p = doc.add_paragraph(style="List Bullet")
                    _docx_add_runs(p, item)
            else:
                p = doc.add_paragraph()
                _docx_add_runs(p, content)

    doc.add_heading("Missing Values", level=1)
    if missing_df is None or missing_df.empty:
        doc.add_paragraph("No missing values were detected in any column.")
    else:
        rows = [[col, f"{int(r['missing_count']):,}", f"{r['missing_pct']}%"] for col, r in missing_df.iterrows()]
        _docx_table(doc, ["Column", "Missing Count", "Missing %"], rows)

    if numeric_summary_df is not None and not numeric_summary_df.empty:
        doc.add_heading("Numeric Column Statistics", level=1)
        cols = list(numeric_summary_df.columns)
        rows = [[col] + [_fmt_num(r[c], is_count=(c == "count")) for c in cols] for col, r in numeric_summary_df.iterrows()]
        _docx_table(doc, ["Column"] + cols, rows)

    if categorical_summary:
        doc.add_heading("Categorical Column Summary", level=1)
        rows = []
        for col, info in categorical_summary.items():
            top = ", ".join(f"{k} ({v})" for k, v in list(info["top_values"].items())[:3])
            rows.append([col, str(info["unique_values"]), top])
        _docx_table(doc, ["Column", "Unique Values", "Top Values"], rows)

    if corr_df is not None and not corr_df.empty:
        doc.add_heading("Correlation Matrix", level=1)
        cols = list(corr_df.columns)
        rows = [[col] + [f"{corr_df.loc[col, c]:.2f}" for c in cols] for col in corr_df.index]
        _docx_table(doc, [""] + cols, rows)

        if chart_fig is not None:
            doc.add_paragraph()
            doc.add_picture(_render_fig_png(chart_fig), width=DocxInches(5.5))

    if outliers:
        doc.add_heading("Potential Outliers (IQR rule)", level=1)
        rows = [[col, str(info["count"]), f"{info['pct']}%"] for col, info in outliers.items()]
        _docx_table(doc, ["Column", "Count", "% of values"], rows)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ==========================================================================
# PPTX - designed to be presented as-is
# ==========================================================================
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5  # 16:9
CHART_MAX_W_IN = 6.5   # consistent "good size, not huge" width for any single chart
GRID_CELL_W_IN = 5.7   # consistent width for charts shown 2-up in a grid


def _pptx_accent_bar(slide, prs, y_in=0, height_in=1.15):
    bar = slide.shapes.add_shape(1, PptxInches(0), PptxInches(y_in), PptxInches(SLIDE_W_IN), PptxInches(height_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_PPTX
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _pptx_title_text(slide, text, top_in, left_in=0.6, width_in=SLIDE_W_IN - 1.2,
                      size=28, color=PptxRGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptxInches(left_in), PptxInches(top_in), PptxInches(width_in), PptxInches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptxPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _pptx_footer(slide, dataset_title, page_no):
    box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(SLIDE_H_IN - 0.45),
                                    PptxInches(SLIDE_W_IN - 1.2), PptxInches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"{dataset_title}  ·  EDA Report"
    run.font.size = PptxPt(9)
    run.font.color.rgb = PptxRGBColor(0xA0, 0xAE, 0xC0)

    box2 = slide.shapes.add_textbox(PptxInches(SLIDE_W_IN - 1.6), PptxInches(SLIDE_H_IN - 0.45),
                                     PptxInches(1.0), PptxInches(0.35))
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = str(page_no)
    run2.font.size = PptxPt(9)
    run2.font.color.rgb = PptxRGBColor(0xA0, 0xAE, 0xC0)


def _pptx_content_slide(prs, heading):
    """A blank slide with the standard accent bar + title, ready for content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_accent_bar(slide, prs)
    _pptx_title_text(slide, heading, top_in=0.28, size=26)
    return slide


def _pptx_bullet_slide(prs, heading, items):
    slide = _pptx_content_slide(prs, heading or "Summary")
    box = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(1.55), PptxInches(SLIDE_W_IN - 1.4), PptxInches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items[:8]:  # cap so slides stay readable
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = PptxPt(12)
        bullet_run = p.add_run()
        bullet_run.text = "●  "
        bullet_run.font.size = PptxPt(18)
        bullet_run.font.color.rgb = ACCENT_PPTX
        bullet_run.font.bold = True
        for segment, is_bold in _split_bold_segments(item):
            run = p.add_run()
            run.text = segment
            run.font.size = PptxPt(18)
            run.font.bold = is_bold
            run.font.color.rgb = PptxRGBColor(0x33, 0x33, 0x33)
    return slide


def _pptx_single_chart_slide(prs, heading, fig, max_width_in=CHART_MAX_W_IN):
    """One chart, centered, capped to a moderate consistent width."""
    slide = _pptx_content_slide(prs, heading)
    left = (SLIDE_W_IN - max_width_in) / 2
    slide.shapes.add_picture(_render_fig_png(fig), PptxInches(left), PptxInches(1.55), width=PptxInches(max_width_in))
    return slide


def _pptx_chart_grid_slides(prs, heading, figs_dict, max_per_slide=4):
    """Multiple charts (e.g. histograms), always the SAME size, laid out
    2-per-row so nothing looks oversized or inconsistent slide to slide.

    eda_utils renders every histogram/bar chart at figsize=(5, 3.2), i.e. a
    fixed 5:3.2 width:height ratio. python-pptx only lets us pin width and
    auto-scales height to match that ratio, so the row height has to be
    derived from the width we pick - not assumed - or rows silently
    overflow the slide.
    """
    if not figs_dict:
        return
    items = list(figs_dict.items())
    fig_aspect = 5 / 3.2  # width / height, matches eda_utils figsize=(5, 3.2)

    top0 = 1.55
    footer_clearance = 0.65
    row_gap, col_gap = 0.3, 0.3
    available_h = SLIDE_H_IN - footer_clearance - top0
    row_h = (available_h - row_gap) / 2       # 2 rows per slide
    cell_w = row_h * fig_aspect
    left0 = (SLIDE_W_IN - (2 * cell_w + col_gap)) / 2

    for chunk_start in range(0, len(items), max_per_slide):
        chunk = items[chunk_start:chunk_start + max_per_slide]
        slide = _pptx_content_slide(prs, heading)
        for i, (_, fig) in enumerate(chunk):
            row, col = divmod(i, 2)
            left = left0 + col * (cell_w + col_gap)
            top = top0 + row * (row_h + row_gap)
            slide.shapes.add_picture(_render_fig_png(fig), PptxInches(left), PptxInches(top),
                                      width=PptxInches(cell_w))


def build_pptx_report(dataset_title, generated_at, model_name, narrative_md,
                       overview, chart_fig=None, hist_figs=None, bar_figs=None) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_H_IN * 914400))
    blank = prs.slide_layouts[6]

    # --- Title slide (no generated-date/model line, per request) ----------
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(SLIDE_W_IN), PptxInches(SLIDE_H_IN))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_PPTX
    bg.line.fill.background()
    bg.shadow.inherit = False

    _pptx_title_text(slide, dataset_title, top_in=2.9, size=40, align=PP_ALIGN.CENTER, left_in=0.8, width_in=SLIDE_W_IN - 1.6)
    _pptx_title_text(slide, "Exploratory Data Analysis Report", top_in=3.8, size=22, bold=False,
                      color=PptxRGBColor(0xD9, 0xE6, 0xF2), align=PP_ALIGN.CENTER, left_in=0.8, width_in=SLIDE_W_IN - 1.6)

    # --- Key statistics slide ------------------------------------------
    slide = _pptx_content_slide(prs, "Key Statistics")
    rows = _overview_rows(overview)
    rows_n, cols_n = len(rows) + 1, 2
    table_shape = slide.shapes.add_table(rows_n, cols_n, PptxInches(2.5), PptxInches(1.8),
                                          PptxInches(SLIDE_W_IN - 5), PptxInches(0.6 * rows_n))
    table = table_shape.table
    table.cell(0, 0).text, table.cell(0, 1).text = "Metric", "Value"
    for c in (0, 1):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_PPTX
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = PptxRGBColor(0xFF, 0xFF, 0xFF)
                r.font.bold = True
    for i, (k, v) in enumerate(rows, start=1):
        table.cell(i, 0).text, table.cell(i, 1).text = k, v
    # --- One slide per narrative section --------------------------------
    for heading, lines in _parse_markdown_sections(narrative_md):
        blocks = _blocks_from_lines(lines)
        items = []
        for kind, content in blocks:
            if kind == "bullets":
                items.extend(content)
            else:
                items.append(content)
        if items:
            _pptx_bullet_slide(prs, heading, items)

    # --- Correlation heatmap: one chart, moderate consistent size --------
    if chart_fig is not None:
        _pptx_single_chart_slide(prs, "Correlation Heatmap", chart_fig)

    # --- Numeric distributions: same-size grid, not oversized -----------
    if hist_figs:
        _pptx_chart_grid_slides(prs, "Numeric Distributions", hist_figs, max_per_slide=4)

    # --- Categorical breakdowns: same-size grid --------------------------
    if bar_figs:
        _pptx_chart_grid_slides(prs, "Categorical Breakdown", bar_figs, max_per_slide=4)

    # --- Closing slide -----------------------------------------------------
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(SLIDE_W_IN), PptxInches(SLIDE_H_IN))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_PPTX
    bg.line.fill.background()
    bg.shadow.inherit = False
    _pptx_title_text(slide, "Thank You", top_in=3.0, size=36, align=PP_ALIGN.CENTER, left_in=0.8, width_in=SLIDE_W_IN - 1.6)
    _pptx_title_text(slide, f"{dataset_title} — EDA Report", top_in=3.9, size=16, bold=False,
                      color=PptxRGBColor(0xD9, 0xE6, 0xF2), align=PP_ALIGN.CENTER, left_in=0.8, width_in=SLIDE_W_IN - 1.6)

    # --- Footers on every content slide (skip title=0 and closing=last) --
    for i, slide in enumerate(prs.slides):
        if i == 0 or i == len(prs.slides) - 1:
            continue
        _pptx_footer(slide, dataset_title, i)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
