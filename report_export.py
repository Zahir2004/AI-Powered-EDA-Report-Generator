"""
report_export.py
-----------------
Turns the AI narrative (Markdown from Groq) + computed statistics into
DOCX, PDF, and PPTX files, each returned as raw bytes ready to hand to
st.download_button. All three read the same input, so the content is
consistent across formats - only the layout differs.

Kept dependency-light on purpose:
- DOCX: python-docx
- PDF:  reportlab (pure Python, no system binaries needed - unlike
        weasyprint/wkhtmltopdf, this installs cleanly on Streamlit Cloud)
- PPTX: python-pptx
"""

import io
import re

from docx import Document
from docx.shared import Pt, Inches as DocxInches

from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
    ListFlowable, ListItem, Preformatted,
)

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt


# --------------------------------------------------------------------------
# Shared markdown parsing (turns Groq's "## Heading" / "- bullet" style
# Markdown into a simple structure every format can render its own way)
# --------------------------------------------------------------------------
def _parse_markdown_sections(markdown_text: str):
    """Split into [(heading_or_None, [raw_lines]), ...] on '## ' headings."""
    sections = []
    heading, lines = None, []
    for line in markdown_text.splitlines():
        if line.startswith("## "):
            if lines or heading:
                sections.append((heading, lines))
            heading, lines = line[3:].strip(), []
        else:
            lines.append(line)
    if lines or heading:
        sections.append((heading, lines))
    return sections


def _blocks_from_lines(lines):
    """Within a section, group into ('bullets', [items]) / ('para', text) blocks."""
    blocks, bullet_buf, para_buf = [], [], []

    def flush():
        if bullet_buf:
            blocks.append(("bullets", bullet_buf[:]))
            bullet_buf.clear()
        if para_buf:
            blocks.append(("para", " ".join(para_buf)))
            para_buf.clear()

    for line in lines:
        stripped = line.strip()
        if not stripped:
            flush()
        elif stripped.startswith(("- ", "* ")):
            if para_buf:
                blocks.append(("para", " ".join(para_buf)))
                para_buf.clear()
            bullet_buf.append(stripped[2:].strip())
        else:
            if bullet_buf:
                blocks.append(("bullets", bullet_buf[:]))
                bullet_buf.clear()
            para_buf.append(stripped)
    flush()
    return blocks


def _strip_markdown_emphasis(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text


# --------------------------------------------------------------------------
# DOCX
# --------------------------------------------------------------------------
def build_docx_report(title, generated_at, model_name, narrative_md, stats_text, chart_fig=None) -> bytes:
    doc = Document()

    doc.add_heading(title, level=0)
    meta = doc.add_paragraph()
    meta.add_run(f"Generated {generated_at} using {model_name}").italic = True

    for heading, lines in _parse_markdown_sections(narrative_md):
        if heading:
            doc.add_heading(heading, level=1)
        for kind, content in _blocks_from_lines(lines):
            if kind == "bullets":
                for item in content:
                    doc.add_paragraph(_strip_markdown_emphasis(item), style="List Bullet")
            else:
                doc.add_paragraph(_strip_markdown_emphasis(content))

    if chart_fig is not None:
        doc.add_heading("Correlation Heatmap", level=1)
        img_buf = io.BytesIO()
        chart_fig.savefig(img_buf, format="png", dpi=150, bbox_inches="tight")
        img_buf.seek(0)
        doc.add_picture(img_buf, width=DocxInches(6))

    doc.add_heading("Computed Statistics", level=1)
    run = doc.add_paragraph().add_run(stats_text)
    run.font.name = "Courier New"
    run.font.size = Pt(8)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------
# PDF (reportlab)
# --------------------------------------------------------------------------
def build_pdf_report(title, generated_at, model_name, narrative_md, stats_text, chart_fig=None) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    meta_style = ParagraphStyle("meta", parent=styles["Normal"], textColor=colors.grey, spaceAfter=14)

    story = [Paragraph(title, styles["Title"]),
             Paragraph(f"Generated {generated_at} using {model_name}", meta_style)]

    for heading, lines in _parse_markdown_sections(narrative_md):
        if heading:
            story.append(Paragraph(heading, styles["Heading2"]))
        for kind, content in _blocks_from_lines(lines):
            if kind == "bullets":
                items = [ListItem(Paragraph(_strip_markdown_emphasis(i), styles["Normal"])) for i in content]
                story.append(ListFlowable(items, bulletType="bullet", leftIndent=18))
            else:
                story.append(Paragraph(_strip_markdown_emphasis(content), styles["Normal"]))
        story.append(Spacer(1, 10))

    if chart_fig is not None:
        img_buf = io.BytesIO()
        chart_fig.savefig(img_buf, format="png", dpi=150, bbox_inches="tight")
        img_buf.seek(0)
        story.append(Paragraph("Correlation Heatmap", styles["Heading2"]))
        story.append(RLImage(img_buf, width=6 * inch, height=4.2 * inch, kind="proportional"))
        story.append(Spacer(1, 10))

    story.append(Paragraph("Computed Statistics", styles["Heading2"]))
    code_style = ParagraphStyle("code", parent=styles["Code"], fontSize=6.5, leading=8)
    story.append(Preformatted(stats_text, code_style))

    doc.build(story)
    return buf.getvalue()


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------
def build_pptx_report(title, generated_at, model_name, narrative_md, chart_fig=None) -> bytes:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Generated {generated_at} using {model_name}"

    # One slide per narrative section
    for heading, lines in _parse_markdown_sections(narrative_md):
        blocks = _blocks_from_lines(lines)
        bullet_items = []
        for kind, content in blocks:
            if kind == "bullets":
                bullet_items.extend(content)
            else:
                # wrap long paragraphs onto the slide as a single bullet
                bullet_items.append(content)
        if not bullet_items:
            continue

        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = heading or "Summary"
        body = slide.placeholders[1].text_frame
        body.clear()
        body.text = _strip_markdown_emphasis(bullet_items[0])[:400]
        for item in bullet_items[1:8]:  # cap to keep slides readable
            p = body.add_paragraph()
            p.text = _strip_markdown_emphasis(item)[:400]
            p.level = 0

    # Optional chart slide
    if chart_fig is not None:
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.6))
        tx.text_frame.text = "Correlation Heatmap"
        tx.text_frame.paragraphs[0].font.size = PptxPt(28)
        tx.text_frame.paragraphs[0].font.bold = True

        img_buf = io.BytesIO()
        chart_fig.savefig(img_buf, format="png", dpi=150, bbox_inches="tight")
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, PptxInches(1), PptxInches(1.1), width=PptxInches(8))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
